"""
pptx_server.py — FastAPI backend for LLM-generated python-pptx execution.

Flow:
  1. POST /generate-pptx  { slides_json, azure_config }
  2. Server calls LLM → asks it to write a self-contained python-pptx script
  3. Script is written to a secure temp file and executed in a subprocess
  4. The resulting .pptx is streamed back to the browser
  5. Both temp files are deleted

Run:
  pip install fastapi uvicorn python-pptx openai
  python pptx_server.py
  # → http://localhost:8765
"""

from __future__ import annotations

import os
import sys
import ast
import json
import asyncio
import tempfile
import subprocess
import textwrap
import traceback
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel

# ─── Optional: httpx for Anthropic direct calls ──────────────────────────────
try:
    import httpx
    HTTPX_AVAILABLE = True
except ImportError:
    HTTPX_AVAILABLE = False

app = FastAPI(title="PPTX Python Server", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:4173", "http://localhost:3000"],
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)

# ─── Request / Response models ────────────────────────────────────────────────
class AzureConfigPayload(BaseModel):
    endpoint: str
    apiKey: str
    deploymentName: str
    apiVersion: Optional[str] = None
    maxTokens: int = 4096
    temperature: float = 0.3


class GeneratePptxRequest(BaseModel):
    documentJson: dict          # full DocumentOutput JSON from the frontend
    azureConfig: AzureConfigPayload
    filename: str = "presentation"


# ─── Direct pptx_elements renderer (no LLM needed) ───────────────────────────
def render_pptx_elements(prs: "Presentation", doc: dict) -> None:
    """Translate pptx_elements arrays directly into python-pptx — pixel-perfect match.
    Coordinate space: 10×7.5 inches (same as PptxGenJS default).
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    SRC_W   = 10.0   # pptx_elements coordinate space (PptxGenJS default)
    SRC_H   = 7.5
    SLIDE_W = 13.33  # output widescreen width
    SLIDE_H = 7.5    # output height (unchanged)
    SCALE_X = SLIDE_W / SRC_W   # 1.333 — stretch x/w to fill widescreen
    SCALE_Y = 1.0               # no vertical change

    ALIGN_MAP = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    VALIGN_MAP = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }

    def to_rgb(hex_str: str) -> RGBColor:
        h = str(hex_str).lstrip("#").strip()
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        if len(h) != 6:
            h = "000000"
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def dim_x(val: object) -> float:
        """x / w: src 10" space → scaled to 13.33" widescreen."""
        if isinstance(val, str) and val.endswith("%"):
            return float(val[:-1]) / 100.0 * SLIDE_W
        return float(val) * SCALE_X  # type: ignore[arg-type]

    def dim_y(val: object) -> float:
        """y / h: src 7.5" space → output 7.5" (no scaling needed)."""
        if isinstance(val, str) and val.endswith("%"):
            return float(val[:-1]) / 100.0 * SLIDE_H
        return float(val) * SCALE_Y  # type: ignore[arg-type]

    for slide_data in doc.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        elements: list = slide_data.get("pptx_elements") or []

        # ── Set background from first full-coverage rect (if any) ──────────
        bg_color = "000000"
        for el in elements:
            if (
                el.get("type") == "rect"
                and dim_x(el.get("w", 0)) >= SLIDE_W * 0.95
                and dim_y(el.get("h", 0)) >= SLIDE_H * 0.90
                and el.get("fill")
            ):
                bg_color = str(el["fill"]).lstrip("#")
                break
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = to_rgb(bg_color)

        # ── Render each element ─────────────────────────────────────────────
        for el in elements:
            el_type = el.get("type", "")
            x = Inches(dim_x(el.get("x", 0)))
            y = Inches(dim_y(el.get("y", 0)))
            w = Inches(max(0.01, dim_x(el.get("w", 1))))
            h = Inches(max(0.01, dim_y(el.get("h", 0.5))))
            fill_hex = str(el.get("fill", "")).lstrip("#")

            if el_type in ("rect", "ellipse"):
                # 1 = MSO_AUTO_SHAPE_TYPE.RECTANGLE, 9 = OVAL
                shape_id = 1 if el_type == "rect" else 9
                shape = slide.shapes.add_shape(shape_id, x, y, w, h)
                if fill_hex:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = to_rgb(fill_hex)
                    # Make the border invisible by matching fill color, zero width
                    shape.line.color.rgb = to_rgb(fill_hex)
                    shape.line.width = Pt(0)
                else:
                    shape.fill.background()
                    shape.line.width = Pt(0)

            elif el_type == "line":
                line_color = fill_hex or str(el.get("color", "CCCCCC")).lstrip("#")
                shape = slide.shapes.add_shape(1, x, y, w, h)  # thin rect as line
                shape.fill.solid()
                shape.fill.fore_color.rgb = to_rgb(line_color)
                shape.line.color.rgb = to_rgb(line_color)
                shape.line.width = Pt(0)

            elif el_type == "text" and el.get("text") is not None:
                txBox = slide.shapes.add_textbox(x, y, w, h)
                tf = txBox.text_frame
                tf.word_wrap = el.get("wrap", True)
                tf.vertical_anchor = VALIGN_MAP.get(
                    str(el.get("valign", "top")), MSO_ANCHOR.TOP
                )
                if fill_hex:
                    txBox.fill.solid()
                    txBox.fill.fore_color.rgb = to_rgb(fill_hex)

                text_val = str(el.get("text", ""))
                lines = text_val.split("\n")
                for i, line_text in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.alignment = ALIGN_MAP.get(
                        str(el.get("align", "left")), PP_ALIGN.LEFT
                    )
                    run = p.add_run()
                    run.text = line_text
                    run.font.size = Pt(float(el.get("fontSize", 14)))
                    run.font.bold = bool(el.get("bold", False))
                    run.font.italic = bool(el.get("italic", False))
                    run.font.color.rgb = to_rgb(el.get("color", "000000"))
                    if el.get("fontFace"):
                        run.font.name = str(el["fontFace"])

        # ── Speaker notes ───────────────────────────────────────────────────
        notes = slide_data.get("speaker_notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)


# ─── Per-slide pptx_elements renderer (used by the main endpoint) ─────────────
def _luminance(hex_str: str) -> float:
    """Return perceived luminance 0.0 (black) → 1.0 (white) for a hex color."""
    h = str(hex_str).lstrip("#").strip()
    if len(h) != 6: return 0.5
    r, g, b = int(h[0:2],16)/255, int(h[2:4],16)/255, int(h[4:6],16)/255
    return 0.2126*r + 0.7152*g + 0.0722*b


def _remap_elements_to_theme(elements: list, theme_colors: dict) -> list:
    """
    Remap arbitrary LLM-chosen colors in pptx_elements to the selected theme palette.
    Only applied when the theme explicitly defines a remap strategy (e.g. zensar_white).
    """
    remap = theme_colors.get("_remap")  # only present for themes that enforce colors
    if not remap:
        return elements

    bg_hex    = theme_colors.get("background", "ffffff")
    header_bg = theme_colors.get("headerBg", "003A70")
    card_fill = theme_colors.get("cardFill", "F9FDFC")
    body_text = theme_colors.get("bodyText", "1A1F2B")
    accent    = theme_colors.get("accent2", "9A1F1F")   # label color
    muted     = theme_colors.get("mutedText", "525A6B")
    positive  = theme_colors.get("positive", "1F6A3A")

    remapped = []
    for el in elements:
        el = dict(el)  # shallow copy so we don't mutate original
        fill = el.get("fill", "")
        color = el.get("color", "")

        if el.get("type") in ("rect", "ellipse") and fill:
            lum = _luminance(fill)
            if lum < 0.12:
                # Very dark fill (near-black) → white slide background
                el["fill"] = bg_hex
            elif lum < 0.35:
                # Dark mid-tone → use theme header/primary color
                el["fill"] = header_bg
            elif lum > 0.88:
                # Near-white fill → card background
                el["fill"] = card_fill

        if el.get("type") == "text" and color:
            lum = _luminance(color)
            if lum > 0.85:
                # Near-white text — keep only if on a dark background (handled by context)
                pass  # leave white text on header bars alone
            elif lum < 0.08:
                # Near-black text → body text color
                el["color"] = body_text
            # Bright saturated colors: detect approximate hue buckets
            else:
                h = str(color).lstrip("#")
                if len(h) == 6:
                    r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
                    # Bright blues → primary navy
                    if b > 150 and b > r*1.5 and b > g*1.2:
                        el["color"] = header_bg
                    # Greens → positive
                    elif g > 150 and g > r*1.3 and g > b*1.3:
                        el["color"] = positive
                    # Reds/oranges → accent label
                    elif r > 150 and r > g*1.4 and r > b*1.4:
                        el["color"] = accent
                    # Mid-grey → muted text
                    elif 80 < r < 180 and abs(r-g) < 30 and abs(g-b) < 30:
                        el["color"] = muted

        remapped.append(el)
    return remapped


def _render_single_slide_elements(slide: "Slide", elements: list, theme_colors: dict) -> None:
    """Render one slide's pptx_elements list onto an already-added blank slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # Remap colors to match selected theme (only for themes with _remap=True, e.g. zensar_white)
    elements = _remap_elements_to_theme(elements, theme_colors)

    SRC_W   = 10.0
    SRC_H   = 7.5
    SLIDE_W = 13.33
    SLIDE_H = 7.5
    SCALE_X = SLIDE_W / SRC_W

    ALIGN_MAP  = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    VALIGN_MAP = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}

    def to_rgb(hex_str: str) -> RGBColor:
        h = str(hex_str).lstrip("#").strip()
        if len(h) == 3: h = "".join(c*2 for c in h)
        if len(h) != 6: h = "000000"
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def dim_x(v: object) -> float:
        if isinstance(v, str) and v.endswith("%"): return float(v[:-1])/100.0 * SLIDE_W
        return float(v) * SCALE_X  # type: ignore[arg-type]

    def dim_y(v: object) -> float:
        if isinstance(v, str) and v.endswith("%"): return float(v[:-1])/100.0 * SLIDE_H
        return float(v)  # type: ignore[arg-type]

    # Set background from first full-coverage rect
    bg_color = theme_colors.get("background", "ffffff")
    for el in elements:
        if el.get("type") == "rect" and dim_x(el.get("w",0)) >= SLIDE_W*0.95 and dim_y(el.get("h",0)) >= SLIDE_H*0.90 and el.get("fill"):
            bg_color = str(el["fill"]).lstrip("#"); break
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = to_rgb(bg_color)

    for el in elements:
        el_type  = el.get("type","")
        x = Inches(dim_x(el.get("x", 0)));  y = Inches(dim_y(el.get("y", 0)))
        w = Inches(max(0.01, dim_x(el.get("w", 1)))); h = Inches(max(0.01, dim_y(el.get("h", 0.5))))
        fill_hex = str(el.get("fill","")).lstrip("#")

        if el_type in ("rect", "ellipse"):
            shape = slide.shapes.add_shape(1 if el_type=="rect" else 9, x, y, w, h)
            if fill_hex:
                shape.fill.solid(); shape.fill.fore_color.rgb = to_rgb(fill_hex)
                shape.line.color.rgb = to_rgb(fill_hex); shape.line.width = Pt(0)
            else:
                shape.fill.background(); shape.line.width = Pt(0)

        elif el_type == "line":
            lc = fill_hex or str(el.get("color","CCCCCC")).lstrip("#")
            s = slide.shapes.add_shape(1, x, y, w, h)
            s.fill.solid(); s.fill.fore_color.rgb = to_rgb(lc)
            s.line.color.rgb = to_rgb(lc); s.line.width = Pt(0)

        elif el_type == "text" and el.get("text") is not None:
            txBox = slide.shapes.add_textbox(x, y, w, h)
            tf = txBox.text_frame; tf.word_wrap = el.get("wrap", True)
            tf.vertical_anchor = VALIGN_MAP.get(str(el.get("valign","top")), MSO_ANCHOR.TOP)
            if fill_hex:
                txBox.fill.solid(); txBox.fill.fore_color.rgb = to_rgb(fill_hex)
            for i, line_text in enumerate(str(el.get("text","")).split("\n")):
                p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
                p.alignment = ALIGN_MAP.get(str(el.get("align","left")), PP_ALIGN.LEFT)
                run = p.add_run(); run.text = line_text
                run.font.size = Pt(float(el.get("fontSize",14)))
                run.font.bold = bool(el.get("bold",False))
                run.font.italic = bool(el.get("italic",False))
                run.font.color.rgb = to_rgb(el.get("color","000000"))
                if el.get("fontFace"): run.font.name = str(el["fontFace"])


def _add_zensar_logo(slide: "Slide") -> None:
    """Stamp the Zensar logo at the standard top-right position on every slide."""
    import os
    from pptx.util import Inches
    logo_path = os.path.join(os.path.dirname(__file__), "public", "zensar_logo.png")
    if not os.path.exists(logo_path):
        return
    # Logo position extracted from actual Zensar template:
    # right edge ~13.13" (11.677 + 1.456), top 0.209" — shift slightly for 13.33" slide
    slide.shapes.add_picture(logo_path, Inches(11.6), Inches(0.18), Inches(1.56), Inches(0.24))


def _add_footer(slide: "Slide", slide_num: int, total: int, title: str) -> None:
    """Add slide number + title footer if not already present in pptx_elements."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    FOOTER_Y = 7.25
    # Title label (left)
    tb_left = slide.shapes.add_textbox(Inches(0.2), Inches(FOOTER_Y), Inches(8.0), Inches(0.22))
    tf = tb_left.text_frame
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = title[:80]
    run.font.size = Pt(8); run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    # Slide N of TOTAL (right-aligned)
    tb_right = slide.shapes.add_textbox(Inches(8.5), Inches(FOOTER_Y), Inches(1.5), Inches(0.22))
    tf2 = tb_right.text_frame
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT; run2 = p2.add_run()
    run2.text = f"Slide {slide_num} of {total}"
    run2.font.size = Pt(8); run2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)


def render_structured_slide(slide: "Slide", slide_data: dict, theme_colors: dict) -> None:
    """Render a slide from its structured fields (title, bullets, table, etc.)
    using a clean, consistent layout. Used when pptx_elements are unavailable."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    def rgb(hex_str: str) -> RGBColor:
        h = str(hex_str).lstrip("#")
        if len(h) == 3: h = "".join(c*2 for c in h)
        if len(h) != 6: h = "334155"
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    HEADER_BG  = theme_colors.get("headerBg",  "1e293b")
    ACCENT     = theme_colors.get("accent",     "3b82f6")
    BODY_TEXT  = theme_colors.get("bodyText",   "1e293b")
    BG         = theme_colors.get("background", "ffffff")
    HEADER_TXT = theme_colors.get("headerText", "ffffff")

    title   = str(slide_data.get("title", ""))
    bullets = slide_data.get("bullets") or []
    subtitle= slide_data.get("subtitle", "")
    layout  = str(slide_data.get("layout", "bullets")).lower()

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(BG)

    # ── Header bar ──────────────────────────────────────────────────────────
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.72))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(HEADER_BG)
    bar.line.width = Pt(0)

    accent_line = slide.shapes.add_shape(1, Inches(0), Inches(0.72), Inches(13.33), Inches(0.06))
    accent_line.fill.solid(); accent_line.fill.fore_color.rgb = rgb(ACCENT)
    accent_line.line.width = Pt(0)

    # Title text in header
    if title:
        txb = slide.shapes.add_textbox(Inches(0.3), Inches(0), Inches(12.7), Inches(0.72))
        tf = txb.text_frame; tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = title
        run.font.size = Pt(22); run.font.bold = True
        run.font.color.rgb = rgb(HEADER_TXT)

    CONTENT_Y = 0.85
    CONTENT_H = 6.4

    if layout in ("title_slide", "title"):
        # Large centered title
        txb = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.1), Inches(3.5))
        tf = txb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = title
        run.font.size = Pt(44); run.font.bold = True; run.font.color.rgb = rgb(HEADER_BG)
        if subtitle:
            txb2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.5))
            tf2 = txb2.text_frame; tf2.word_wrap = True
            p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run(); run2.text = subtitle
            run2.font.size = Pt(22); run2.font.color.rgb = rgb(BODY_TEXT)

    elif layout == "two_column":
        left  = slide_data.get("left_column") or []
        right = slide_data.get("right_column") or []
        lt    = slide_data.get("left_title", "")
        rt    = slide_data.get("right_title", "")
        for col_x, col_items, col_title, col_accent in [
            (0.3, left, lt, ACCENT), (6.87, right, rt, theme_colors.get("accent2", ACCENT))
        ]:
            card = slide.shapes.add_shape(1, Inches(col_x), Inches(CONTENT_Y), Inches(6.2), Inches(CONTENT_H))
            card.fill.solid(); card.fill.fore_color.rgb = rgb("f8fafc")
            card.line.color.rgb = rgb("e2e8f0"); card.line.width = Pt(0.5)
            top = slide.shapes.add_shape(1, Inches(col_x), Inches(CONTENT_Y), Inches(6.2), Inches(0.06))
            top.fill.solid(); top.fill.fore_color.rgb = rgb(col_accent); top.line.width = Pt(0)
            ty = CONTENT_Y + 0.1
            if col_title:
                txb = slide.shapes.add_textbox(Inches(col_x+0.15), Inches(ty), Inches(5.9), Inches(0.4))
                p = txb.text_frame.paragraphs[0]; run = p.add_run()
                run.text = str(col_title); run.font.bold = True; run.font.size = Pt(14)
                run.font.color.rgb = rgb(col_accent); ty += 0.45
            if col_items:
                txb = slide.shapes.add_textbox(Inches(col_x+0.15), Inches(ty), Inches(5.9), Inches(CONTENT_H-0.6))
                tf = txb.text_frame; tf.word_wrap = True
                for i, item in enumerate(col_items):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.space_before = Pt(4)
                    run = p.add_run(); run.text = f"• {item}"
                    run.font.size = Pt(14); run.font.color.rgb = rgb(BODY_TEXT)

    elif layout == "table" and slide_data.get("table"):
        tbl = slide_data["table"]
        headers = tbl.get("headers", [])
        rows    = tbl.get("rows", [])
        if headers and rows:
            col_count = len(headers)
            row_count = len(rows) + 1
            tbl_shape = slide.shapes.add_table(row_count, col_count,
                Inches(0.3), Inches(CONTENT_Y), Inches(12.7), Inches(min(CONTENT_H, row_count * 0.45))).table
            for ci, h in enumerate(headers):
                cell = tbl_shape.cell(0, ci)
                cell.text = str(h)
                cell.fill.solid(); cell.fill.fore_color.rgb = rgb(HEADER_BG)
                run = cell.text_frame.paragraphs[0].runs[0] if cell.text_frame.paragraphs[0].runs else cell.text_frame.paragraphs[0].add_run()
                run.font.bold = True; run.font.color.rgb = rgb("ffffff"); run.font.size = Pt(13)
            for ri, row in enumerate(rows):
                for ci, cell_val in enumerate(row):
                    cell = tbl_shape.cell(ri+1, ci)
                    cell.text = str(cell_val)
                    if ri % 2 == 0:
                        cell.fill.solid(); cell.fill.fore_color.rgb = rgb("f1f5f9")
                    run = cell.text_frame.paragraphs[0].runs[0] if cell.text_frame.paragraphs[0].runs else cell.text_frame.paragraphs[0].add_run()
                    run.font.size = Pt(12); run.font.color.rgb = rgb(BODY_TEXT)

    else:
        # Default: bullets list
        items = bullets or ([subtitle] if subtitle else [])
        if items:
            txb = slide.shapes.add_textbox(Inches(0.5), Inches(CONTENT_Y + 0.15), Inches(12.3), Inches(CONTENT_H))
            tf = txb.text_frame; tf.word_wrap = True
            row_h = min(0.75, (CONTENT_H - 0.2) / max(len(items), 1))
            fs = 18 if row_h >= 0.6 else (15 if row_h >= 0.45 else 13)
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(max(2, int(row_h * 8)))
                run = p.add_run(); run.text = f"• {item}"
                run.font.size = Pt(fs); run.font.color.rgb = rgb(BODY_TEXT)

    # Speaker notes
    notes = slide_data.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)


# ─── Theme color lookup ────────────────────────────────────────────────────────
THEME_COLORS: dict[str, dict] = {
    "corporate_blue":  {"headerBg":"1e3a5f","accent":"3b82f6","accent2":"0ea5e9","bodyText":"1e293b","background":"ffffff","headerText":"ffffff"},
    "dark_professional":{"headerBg":"0f172a","accent":"6366f1","accent2":"8b5cf6","bodyText":"1e293b","background":"f8fafc","headerText":"ffffff"},
    "green_growth":    {"headerBg":"14532d","accent":"22c55e","accent2":"16a34a","bodyText":"1e293b","background":"ffffff","headerText":"ffffff"},
    "red_bold":        {"headerBg":"7f1d1d","accent":"ef4444","accent2":"f97316","bodyText":"1e293b","background":"ffffff","headerText":"ffffff"},
    "purple_creative": {"headerBg":"3b0764","accent":"a855f7","accent2":"ec4899","bodyText":"1e293b","background":"ffffff","headerText":"ffffff"},
    # Zensar brand: white bg, navy #003A70, crimson labels #9A1F1F, steel gray #525A6B, forest green #1F6A3A
    "zensar_white":    {"headerBg":"003A70","accent":"003A70","accent2":"9A1F1F","bodyText":"1A1F2B","background":"ffffff","headerText":"ffffff","cardFill":"F9FDFC","labelColor":"9A1F1F","mutedText":"525A6B","positive":"1F6A3A","_remap":True},
    "default":         {"headerBg":"1e293b","accent":"3b82f6","accent2":"0ea5e9","bodyText":"1e293b","background":"ffffff","headerText":"ffffff"},
}

def get_theme_colors(doc: dict) -> dict:
    theme = str(doc.get("theme", "")).lower().replace(" ", "_")
    return THEME_COLORS.get(theme, THEME_COLORS["default"])



# NOTE: 'os' is intentionally NOT banned — the script needs os.environ["PPTX_OUT_PATH"]
BANNED_IMPORTS = {
    "subprocess", "shutil", "socket", "requests", "urllib",
    "httpx", "http", "ftplib", "smtplib", "paramiko", "fabric",
    "ctypes", "cffi", "importlib", "pickle", "marshal",
    "__import__", "eval", "exec", "compile",
    "builtins", "sys",
}
BANNED_PATTERNS = [
    "subprocess", "os.system", "os.popen", "os.exec", "os.spawn",
    "os.remove", "os.unlink", "os.rmdir", "os.rename", "os.listdir",
    "os.walk", "os.makedirs", "os.mkdir",
    "shutil", "__import__", "importlib",
    "open(", "socket.", "requests.", "urllib.",
    "pickle", "marshal", "eval(", "exec(",
]


def is_safe_code(code: str) -> tuple[bool, str]:
    """Return (safe, reason). Blocks code with dangerous patterns."""
    code_lower = code.lower()
    for pat in BANNED_PATTERNS:
        if pat.lower() in code_lower:
            return False, f"Blocked pattern found: '{pat}'"
    # AST import check
    try:
        tree = ast.parse(code)
        for node in ast.walk(tree):
            if isinstance(node, (ast.Import, ast.ImportFrom)):
                names = (
                    [a.name for a in node.names]
                    if isinstance(node, ast.Import)
                    else [node.module or ""]
                )
                for name in names:
                    root = name.split(".")[0]
                    if root in BANNED_IMPORTS:
                        return False, f"Banned import: '{root}'"
    except SyntaxError:
        pass  # will fail at execution — let subprocess handle it
    return True, "ok"


# ─── Endpoint helpers + LLM call ─────────────────────────────────────────────
def _strip_to_base(endpoint: str) -> str:
    base = endpoint.rstrip("/")
    for suffix in [
        "/anthropic/v1/messages",
        "/models/chat/completions",
        "/openai/v1/chat/completions",
        "/openai/v1",
        "/openai/deployments",
        "/v1",
    ]:
        if base.endswith(suffix):
            base = base[: -len(suffix)].rstrip("/")
            break
    return base


def _is_azure_openai(ep: str) -> bool:
    return ".openai.azure.com" in ep.lower() or "cognitiveservices.azure.com" in ep.lower()


def _is_services(ep: str) -> bool:
    return "services.ai.azure.com" in ep.lower()


def _is_openai_model(deployment: str) -> bool:
    d = deployment.lower()
    return d.startswith("gpt-") or d.startswith("o1") or d.startswith("o3") or d.startswith("o4") or d.startswith("codex")


# ─── System prompt for code generation ────────────────────────────────────────
SYSTEM_PROMPT = textwrap.dedent("""
You are an expert Python developer specialising in python-pptx.

Your task: given a JSON object describing a presentation, write a COMPLETE, 
self-contained Python script that uses ONLY the `pptx` library (python-pptx) 
to generate a beautiful, professional PowerPoint file.

STRICT RULES — violating any of these will break the pipeline:
1. Output ONLY raw Python code. No markdown fences (```python), no explanation.
2. The script MUST write the file to the path stored in the environment variable
   PPTX_OUT_PATH — use:  prs.save(os.environ["PPTX_OUT_PATH"])
3. You MUST import os at the top of the script (it is the ONLY standard-library
   import allowed for reading PPTX_OUT_PATH).
4. Use ONLY these libraries: pptx, os  — nothing else.
5. Do NOT use: subprocess, socket, requests, urllib, open(), eval(), exec(),
   shutil, sys, importlib, or any network/filesystem operation other than
   prs.save(...).
6. The presentation must be widescreen 16:9 (13.33 × 7.5 inches).
7. Every slide must have a coloured header bar with the slide title.
8. Make the design visually rich: use colour fills, accent bars, stat cards,
   clean typography — match the theme extracted from the JSON.
9. Add speaker notes to each slide using slide.notes_slide.notes_text_frame.text
10. Handle all slide layouts: title, bullets, two_column, table, stats, closing,
    section_divider, agenda, quote.
11. The script must not crash — handle missing or empty fields gracefully with
    sensible defaults.
""").strip()


def build_user_prompt(doc: dict) -> str:
    slim = json.loads(json.dumps(doc))
    for slide in slim.get("slides", []):
        slide.pop("html", None)
        slide.pop("background_html", None)
        # Keep pptx_elements so LLM can see the exact design
    return (
        "Generate a complete python-pptx script for the following presentation JSON.\n"
        "IMPORTANT: Each slide has a 'pptx_elements' array containing the EXACT shapes, "
        "colors, positions, and text of the original design (coordinate space: 10×7.5 inches). "
        "Reproduce these elements faithfully using python-pptx — same colors (hex), same positions, "
        "same fonts and sizes. Do NOT invent a new design.\n\n"
        f"```json\n{json.dumps(slim, indent=2)}\n```"
    )


# ─── Core LLM caller (shared by single-doc and per-slide paths) ──────────────
def _call_llm_messages(system: str, user_msg: str, cfg: AzureConfigPayload,
                       max_tokens: Optional[int] = None) -> str:
    """Route to Anthropic or OpenAI, return stripped code string."""
    tok = max_tokens or cfg.maxTokens

    if _is_services(cfg.endpoint) and _is_claude(cfg.deploymentName):
        if not HTTPX_AVAILABLE:
            raise RuntimeError("httpx not installed. Run: pip install httpx")
        base = _strip_to_base(cfg.endpoint)
        resp = httpx.post(
            f"{base}/anthropic/v1/messages",
            headers={"Content-Type": "application/json",
                     "x-api-key": cfg.apiKey,
                     "anthropic-version": "2023-06-01"},
            json={"model": cfg.deploymentName, "max_tokens": tok,
                  "temperature": cfg.temperature, "system": system,
                  "messages": [{"role": "user", "content": user_msg}]},
            timeout=180,
        )
        if not resp.is_success:
            raise RuntimeError(f"Error code: {resp.status_code} - {resp.text}")
        code = resp.json()["content"][0]["text"]
    else:
        try:
            from openai import OpenAI
        except ImportError:
            raise RuntimeError("openai package not installed. Run: pip install openai")
        base = _strip_to_base(cfg.endpoint)
        if _is_azure_openai(cfg.endpoint):
            base_url = f"{base}/openai/deployments/{cfg.deploymentName}"
        elif _is_services(cfg.endpoint) and _is_openai_model(cfg.deploymentName):
            # Foundry unified endpoint + GPT/o-series → Azure OpenAI path
            base_url = f"{base}/openai/deployments/{cfg.deploymentName}"
        elif _is_services(cfg.endpoint):
            # Foundry unified endpoint + non-OpenAI (Llama, Phi…) → AI Inference path
            base_url = f"{base}/models"
        else:
            base_url = f"{base}/v1"
        api_version = cfg.apiVersion or ("2025-04-01-preview" if _is_services(cfg.endpoint) else None)
        kwargs: dict = {"api_key": cfg.apiKey, "base_url": base_url,
                        "default_headers": {"api-key": cfg.apiKey}}
        if api_version:
            kwargs["default_query"] = {"api-version": api_version}
        client = OpenAI(**kwargs)
        response = client.chat.completions.create(
            model=cfg.deploymentName,
            messages=[{"role": "system", "content": system},
                      {"role": "user", "content": user_msg}],
            max_tokens=tok,
            temperature=cfg.temperature,
            timeout=180,
        )
        code = response.choices[0].message.content or ""

    if "```" in code:
        code = "\n".join(l for l in code.splitlines() if not l.strip().startswith("```"))
    return code.strip()


def call_llm(doc: dict, cfg: AzureConfigPayload) -> str:
    """Single full-doc LLM call (kept for /preview-code endpoint)."""
    return _call_llm_messages(SYSTEM_PROMPT, build_user_prompt(doc), cfg)


# ─── Per-slide prompt & parallel LLM ─────────────────────────────────────────
PER_SLIDE_SYSTEM = textwrap.dedent("""
    You are a python-pptx expert. Write ONLY a Python function for exactly ONE slide.

    FUNCTION SIGNATURE: def add_slide(prs):

    STRICT RULES — any violation breaks the pipeline:
    1. First line of function body: slide = prs.slides.add_slide(prs.slide_layouts[6])
    2. Slide is 13.33" wide × 7.5" tall (already set on prs — do NOT change it).
    3. These names are already in scope (imported at module level):
       Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR
    4. pptx_elements coordinate space is 10" wide → multiply all x/w values by 1.333.
       Color hex strings have NO # prefix: RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
       rect → add_shape(1,...), ellipse → add_shape(9,...), text → add_textbox(...)
    5. DO NOT add import statements inside the function.
    6. DO NOT call prs.save().
    7. DO NOT use open(), subprocess, os.system, eval, exec, or network calls.
    8. DO NOT output markdown code fences — raw Python only.
    9. Add speaker_notes if present: slide.notes_slide.notes_text_frame.text = "..."
    10. If no pptx_elements, use title/bullets/layout fields to build a clean, professional slide.
""").strip()


async def _generate_slide_func_llm(
    slide_data: dict, slide_num: int, total: int, cfg: AzureConfigPayload
) -> str:
    """Call LLM in a thread to write python-pptx code for one slide."""
    slim = {k: v for k, v in slide_data.items() if k not in ("html", "background_html")}
    user_msg = (
        f"Slide {slide_num} of {total}:\n\n"
        f"```json\n{json.dumps(slim, indent=2)}\n```\n\n"
        "Write `def add_slide(prs):` for this slide."
    )
    # Per-slide token cap: 4096 is enough for one slide's code
    return await asyncio.to_thread(_call_llm_messages, PER_SLIDE_SYSTEM, user_msg, cfg, 4096)


def _build_combined_script(slide_funcs: list, doc: dict) -> str:
    """Combine per-slide function strings into one executable python-pptx script."""
    header = textwrap.dedent("""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        import os
    """).strip()

    funcs = []
    calls = []
    for i, func_code in enumerate(slide_funcs):
        # Rename def add_slide( → def add_slide_N( to avoid name collisions
        renamed = func_code.replace("def add_slide(", f"def add_slide_{i + 1}(", 1)
        funcs.append(renamed)
        calls.append(f"add_slide_{i + 1}(prs)")

    title = str(doc.get("title", "Presentation")).replace('"', '\\"')
    main = textwrap.dedent(f"""
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        prs.core_properties.title = "{title}"
        {chr(10).join(calls)}
        prs.save(os.environ["PPTX_OUT_PATH"])
    """).strip()

    return header + "\n\n" + "\n\n".join(funcs) + "\n\n" + main


# ─── Execute in subprocess ────────────────────────────────────────────────────
def run_code_in_subprocess(code: str, out_path: str) -> tuple[bool, str]:
    """
    Write code to a temp .py file and execute it with the Python interpreter.
    Returns (success, stderr_or_error_message).
    """
    with tempfile.NamedTemporaryFile(
        mode="w", suffix=".py", delete=False, prefix="pptx_gen_"
    ) as f:
        f.write(code)
        script_path = f.name

    try:
        result = subprocess.run(
            [sys.executable, script_path],
            capture_output=True,
            text=True,
            timeout=60,
            env={
                **os.environ,
                "PPTX_OUT_PATH": out_path,
                # Restrict PATH — no network tools available via env
                "PATH": "/usr/bin:/bin",
            },
        )
        if result.returncode != 0:
            return False, result.stderr or result.stdout or "Unknown error"
        return True, ""
    except subprocess.TimeoutExpired:
        return False, "Script execution timed out (60s limit exceeded)"
    except Exception as e:
        return False, str(e)
    finally:
        Path(script_path).unlink(missing_ok=True)


# ─── Health check ─────────────────────────────────────────────────────────────
@app.get("/health")
def health():
    return {"status": "ok", "httpx_available": HTTPX_AVAILABLE}


# ─── Preview endpoint: returns generated code only (no execution) ─────────────
@app.post("/preview-code")
async def preview_code(req: GeneratePptxRequest):
    """Returns the LLM-generated python-pptx code without executing it."""
    try:
        code = call_llm(req.documentJson, req.azureConfig)
        safe, reason = is_safe_code(code)
        return JSONResponse({
            "code": code,
            "safe": safe,
            "safetyNote": reason,
            "lines": len(code.splitlines()),
        })
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── Main endpoint: generate + execute + return file ─────────────────────────
@app.post("/generate-pptx")
async def generate_pptx(req: GeneratePptxRequest):
    doc = req.documentJson
    slides = doc.get("slides", [])

    if not slides:
        raise HTTPException(status_code=400, detail="No slides in document.")

    with tempfile.NamedTemporaryFile(
        suffix=".pptx", delete=False, prefix="pptx_out_"
    ) as f:
        out_path = f.name

    try:
        # ── Fast path: all slides have pptx_elements → direct render, no LLM ──
        if all(slide.get("pptx_elements") for slide in slides):
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)
            prs.core_properties.title  = doc.get("title", "Presentation")
            prs.core_properties.author = doc.get("author", "ZenSpark")
            theme_colors = get_theme_colors(doc)
            theme_name = str(doc.get("theme", "")).lower().replace(" ", "_")
            total_slides = len(slides)
            doc_title = doc.get("title", "Presentation")
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                elements = slide_data.get("pptx_elements", [])
                _render_single_slide_elements(slide, elements, theme_colors)
                # Auto-inject Zensar logo top-right when using zensar_white theme
                if theme_name == "zensar_white":
                    _add_zensar_logo(slide)
                # Add footer if LLM didn't include "Slide N of" text already
                slide_num = slide_data.get("slide_number", slides.index(slide_data) + 1)
                has_footer = any("of" in str(el.get("text","")).lower() and "slide" in str(el.get("text","")).lower() for el in elements)
                if not has_footer:
                    _add_footer(slide, slide_num, total_slides, doc_title)
                if slide_data.get("speaker_notes"):
                    slide.notes_slide.notes_text_frame.text = str(slide_data["speaker_notes"])
            prs.save(out_path)

        # ── Parallel LLM path: one LLM call per slide, all concurrent ────────
        else:
            tasks = [
                _generate_slide_func_llm(slide_data, i + 1, len(slides), req.azureConfig)
                for i, slide_data in enumerate(slides)
            ]
            # All LLM calls run concurrently — total time ≈ slowest single slide
            slide_funcs = await asyncio.gather(*tasks)

            combined = _build_combined_script(list(slide_funcs), doc)

            safe, reason = is_safe_code(combined)
            if not safe:
                Path(out_path).unlink(missing_ok=True)
                raise HTTPException(
                    status_code=422,
                    detail=f"Generated code failed safety check: {reason}",
                )

            success, err = run_code_in_subprocess(combined, out_path)
            if not success:
                Path(out_path).unlink(missing_ok=True)
                raise HTTPException(
                    status_code=500,
                    detail=f"Script execution failed:\n{err}",
                )

        if not Path(out_path).exists() or Path(out_path).stat().st_size == 0:
            Path(out_path).unlink(missing_ok=True)
            raise HTTPException(status_code=500, detail="Failed to produce a .pptx file.")

        safe_filename = "".join(
            c for c in req.filename if c.isalnum() or c in "-_ "
        ).strip() or "presentation"

        return FileResponse(
            path=out_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{safe_filename}.pptx",
            background=_cleanup(out_path),
        )

    except HTTPException:
        raise
    except Exception as e:
        Path(out_path).unlink(missing_ok=True)
        raise HTTPException(status_code=500, detail=str(e))



# ─── Cleanup helper ───────────────────────────────────────────────────────────
from starlette.background import BackgroundTask

def _cleanup(path: str) -> BackgroundTask:
    def _delete():
        Path(path).unlink(missing_ok=True)
    return BackgroundTask(_delete)


# ─── Dev server entry point ───────────────────────────────────────────────────
if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PPTX_SERVER_PORT", 8765))
    print(f"🚀  PPTX Python Server → http://localhost:{port}")
    print(f"    Endpoints: GET /health · POST /generate-pptx · POST /preview-code")
    uvicorn.run("pptx_server:app", host="0.0.0.0", port=port, reload=True)
