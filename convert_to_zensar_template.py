"""
Convert AA_ZenseAI.QI Offerings V7.0.pptx → Zensar template style
- Same content, same layout intent
- White background, Zensar navy #003A70 headers, crimson #9A1F1F accents
- Calibri font throughout
- Zensar logo top-right on every slide
- Slide footer: "ZenseAI.QI  •  Zensar Technologies  |  Confidential" + slide number
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand palette ──────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x00, 0x3A, 0x70)   # Zensar navy
CRIMSON   = RGBColor(0x9A, 0x1F, 0x1F)   # Zensar crimson
BODY      = RGBColor(0x1A, 0x1F, 0x2B)   # near-black body text
MUTED     = RGBColor(0x52, 0x5A, 0x6B)   # steel gray sub-text
GREEN     = RGBColor(0x1F, 0x6A, 0x3A)   # positive/green
CARD      = RGBColor(0xF9, 0xFD, 0xFC)   # very light card bg
DIVIDER   = RGBColor(0xD0, 0xD9, 0xE6)   # light divider
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # amber badge
LOGO_PATH = os.path.join(os.path.dirname(__file__), "public", "zensar_logo.png")

SLIDE_W = 13.33
SLIDE_H = 7.5


def rgb(hex6: str) -> RGBColor:
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor = WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color: RGBColor, line=False) -> "Shape":
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line:
        shape.line.color.rgb = fill_color
        shape.line.width = Pt(0)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=12, bold=False, italic=False,
             color: RGBColor = BODY,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             wrap=True, font_face="Calibri",
             fill_color=None) -> "Shape":
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    if fill_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill_color
    else:
        txBox.fill.background()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_face
    return txBox


def add_multiline_text(slide, lines, x, y, w, h,
                       font_size=12, bold=False, color: RGBColor = BODY,
                       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                       wrap=True, font_face="Calibri",
                       fill_color=None, line_spacing_pt=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    if fill_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill_color
    else:
        txBox.fill.background()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_face
    return txBox


def add_logo(slide):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(11.6), Inches(0.18), Inches(1.56), Inches(0.24))


def add_footer(slide, slide_num: int, total: int, left_text="ZenseAI.QI  •  Zensar Technologies  |  Confidential"):
    # Thin navy line above footer
    add_rect(slide, 0.2, 7.2, 12.93, 0.02, NAVY)
    add_text(slide, left_text, 0.2, 7.23, 9.0, 0.22,
             font_size=7, color=MUTED, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, f"{slide_num} / {total}", 10.5, 7.23, 2.0, 0.22,
             font_size=7, color=MUTED, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def add_header_bar(slide, title: str, tag: str = "", subtitle: str = ""):
    """Standard Zensar content-slide header: navy top bar with title."""
    add_rect(slide, 0, 0, SLIDE_W, 1.05, NAVY)
    # Thin crimson accent line
    add_rect(slide, 0, 1.05, SLIDE_W, 0.04, CRIMSON)
    if tag:
        add_text(slide, tag, 0.3, 0.06, 3.0, 0.22,
                 font_size=9, bold=True, color=CRIMSON, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, 0.3, 0.25, 10.5, 0.65,
             font_size=24, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 1.12, 10.5, 0.3,
                 font_size=11, bold=False, color=MUTED, valign=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_cover(slide, slide_num, total):
    set_bg(slide, WHITE)
    # Large navy hero block
    add_rect(slide, 0, 0, SLIDE_W, 5.8, NAVY)
    # Crimson accent bar at bottom of hero
    add_rect(slide, 0, 5.8, SLIDE_W, 0.08, CRIMSON)
    # Decorative circle (top right)
    c = slide.shapes.add_shape(9, Inches(9.5), Inches(-1.0), Inches(4.5), Inches(4.5))
    c.fill.solid(); c.fill.fore_color.rgb = rgb("002D5B")
    c.line.fill.background()
    # Tag
    add_text(slide, "ZenseAI.QI  •  PROPOSAL", 0.9, 0.55, 4.0, 0.3,
             font_size=9, bold=True, color=CRIMSON)
    # Zensar logo top right
    add_logo(slide)
    # Main title
    add_text(slide, "AI-Led Quality Intelligence for The AA", 0.7, 1.5, 11.5, 1.5,
             font_size=40, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE, wrap=True)
    add_text(slide, "Proof of Value  |  ZenseAI.QI Platform  |  Zensar Technologies",
             0.7, 3.1, 11.5, 0.5, font_size=16, bold=False, color=rgb("B8C2E0"))
    # Date bar
    add_rect(slide, 0, 5.88, SLIDE_W, 1.62, WHITE)
    add_text(slide, "June 2026", 0.9, 6.1, 3.0, 0.3, font_size=11, bold=True, color=CRIMSON)
    add_footer(slide, slide_num, total)


def build_agenda(slide, slide_num, total):
    set_bg(slide, WHITE)
    add_header_bar(slide, "What we will cover today", tag="OVERVIEW",
                   subtitle="ZenseAI.QI  •  AI-Led Quality Intelligence")
    add_logo(slide)

    items = [
        ("01", "Testing Lifecycle Today at The AA",
         "Current projects, tech stack, automation coverage and AI tools in use across all three delivery towers"),
        ("02", "AA – Key Challenges & Pain Points",
         "Flaky tests, manual document validation, no AI-led test design, and gaps in quality observability"),
        ("03", "Testing Lifecycle Reimagined with AA-AI.QI",
         "How ZenseAI.QI agents transform each stage — from requirements to test design, execution and insights"),
        ("04", "ZenseAI.QI – AI-Led QE Foundry (Demo)",
         "Live demonstration of the ZenseAI.QI platform — AI-powered QE workflows built for teams like AA"),
        ("05", "Proposed POV — Phased Implementation Plan",
         "A 3-phase roadmap starting with a zero-risk PDF validation POC, expanding to full QE intelligence"),
        ("06", "Implementation Prerequisites, Skills & Success Measures",
         "Infra, access, LLM, InfoSec whitelisting requirements and how we measure success"),
    ]

    col_w = 5.8
    cols = [(0.4, 1.55), (6.7, 1.55)]
    rows_per_col = 3

    for idx, (num, title, desc) in enumerate(items):
        col = idx // rows_per_col
        row = idx % rows_per_col
        cx, cy_start = cols[col]
        row_h = 1.6
        y = cy_start + row * row_h

        # Number badge
        add_rect(slide, cx, y + 0.05, 0.38, 0.95, NAVY)
        add_text(slide, num, cx, y + 0.05, 0.38, 0.95,
                 font_size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # Card
        add_rect(slide, cx + 0.42, y, col_w - 0.45, 1.05, CARD)
        add_text(slide, title, cx + 0.52, y + 0.05, col_w - 0.65, 0.3,
                 font_size=12, bold=True, color=NAVY)
        add_text(slide, desc, cx + 0.52, y + 0.35, col_w - 0.65, 0.65,
                 font_size=9.5, color=BODY, wrap=True)

    add_footer(slide, slide_num, total)


def build_pipeline_slide(slide, slide_num, total, tag: str, title: str,
                         cols_data: list, bottom_row: list = None):
    """Generic chevron-pipeline layout used for slides 3, 5, 8."""
    set_bg(slide, WHITE)
    add_header_bar(slide, title, tag=tag)
    add_logo(slide)

    stages = ["Requirements\n& Stories", "Test Design\n(TC Authoring)",
              "Automation\nBuild", "Test\nExecution",
              "Defect Mgmt\n& Triage", "Specialized\n& Reporting"]

    # Chevron row
    chevron_w = 1.95
    chevron_h = 0.6
    chevron_y = 1.45
    for i, label in enumerate(stages):
        cx = 0.35 + i * (chevron_w + 0.07)
        fill = NAVY if i > 0 else CRIMSON
        # draw as rect (true chevrons require custom XML; rect is readable)
        add_rect(slide, cx, chevron_y, chevron_w, chevron_h, fill)
        add_text(slide, f"{i+1:02d}", cx + 0.05, chevron_y, 0.35, chevron_h,
                 font_size=7, bold=True, color=rgb("CBD5E1"),
                 valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_text(slide, label, cx + 0.4, chevron_y, chevron_w - 0.42, chevron_h,
                 font_size=9, bold=True, color=WHITE,
                 valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, wrap=True)

    # Column content
    col_w = 1.87
    for ci, col in enumerate(cols_data):
        cx = 0.35 + ci * (col_w + 0.07)
        agent_name = col.get("agent", "")
        badge = col.get("badge", "")
        badge_color = col.get("badge_color", None)
        points = col.get("points", [])
        poc_target = col.get("poc_target", "")
        later = col.get("later", "")

        y = 2.2
        # Agent name
        add_text(slide, agent_name, cx, y, col_w, 0.28,
                 font_size=11, bold=True, color=NAVY)
        y += 0.28

        # Badge
        if badge:
            bc = rgb(badge_color) if badge_color else GREEN
            add_rect(slide, cx, y, 0.7, 0.22, bc)
            add_text(slide, badge, cx, y, 0.7, 0.22,
                     font_size=7, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            y += 0.26

        # Bullet points
        for pt in points:
            add_text(slide, f"• {pt}", cx + 0.05, y, col_w - 0.05, 0.34,
                     font_size=9.5, color=BODY, wrap=True)
            y += 0.36

        # POC target
        if poc_target:
            y = max(y, 4.6)
            add_text(slide, "POC TARGET", cx, y, col_w, 0.18,
                     font_size=7, bold=True, color=GREEN)
            add_text(slide, poc_target, cx, y + 0.18, col_w, 0.32,
                     font_size=9, bold=True, color=BODY, wrap=True)
            y += 0.54

        if later:
            add_text(slide, "PICK UP LATER", cx, y + 0.05, col_w, 0.18,
                     font_size=7, bold=True, color=AMBER)
            add_text(slide, later, cx, y + 0.23, col_w, 0.32,
                     font_size=9, color=BODY, wrap=True)

    # Bottom row summary
    if bottom_row:
        by = 6.4
        col_w2 = 12.5 / len(bottom_row)
        for bi, (label, desc) in enumerate(bottom_row):
            bx = 0.35 + bi * col_w2
            add_text(slide, label, bx, by, col_w2 - 0.1, 0.25,
                     font_size=10, bold=True, color=NAVY)
            add_text(slide, desc, bx, by + 0.25, col_w2 - 0.1, 0.2,
                     font_size=8.5, color=MUTED)

    add_footer(slide, slide_num, total)


def build_challenges(slide, slide_num, total):
    set_bg(slide, WHITE)
    add_header_bar(slide, "AA – Key Focus Areas for Improvement", tag="THE AA",
                   subtitle="Where we can mitigate the quality risk and reduce manual effort")
    add_logo(slide)

    challenges = [
        ("🐛  Reduction of Flaky Tests",
         "BCAS and Payment Gateway suffer persistent test failures driven by AUT slowness, "
         "timeout issues and corrupt test data. High rerun cost."),
        ("📄  Manual PDF / Document Validation",
         "500+ Paragon/Lace documents (policies, letters, renewals) validated entirely by hand. "
         "No automated content, pricing or alignment checks."),
        ("🔍  AI-Led Test Design",
         "Requirements analysis and test case generation are manual. GitHub Copilot is used only "
         "for code assist — no intelligent test generation from stories."),
        ("📊  Quality Observability",
         "No centralized dashboard for defect trends, test health or pipeline quality signals "
         "across Azure DevOps and Jira/Zephyr."),
        ("🔒  Integrate Security & Performance",
         "Payment Gateway and FHUB APIs have no automated security testing. Performance "
         "regression is manual for high-traffic paths."),
        ("📋  Zephyr Limitations",
         "No test-case-level execution tracking in Zephyr. Poor Jira linking between defects "
         "and test cases. Undocumented business logic."),
    ]

    cols = 3
    card_w = 3.9
    card_h = 2.05
    gap_x = 0.22
    start_x = 0.35
    start_y = 1.55

    for idx, (title, desc) in enumerate(challenges):
        col = idx % cols
        row = idx // cols
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + 0.15)

        # Card bg
        add_rect(slide, cx, cy, card_w, card_h, CARD)
        # Navy top strip
        add_rect(slide, cx, cy, card_w, 0.35, NAVY)
        add_text(slide, title, cx + 0.12, cy + 0.02, card_w - 0.15, 0.32,
                 font_size=11, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, desc, cx + 0.12, cy + 0.4, card_w - 0.2, card_h - 0.48,
                 font_size=10, color=BODY, wrap=True)

    add_footer(slide, slide_num, total)


def build_platform_overview(slide, slide_num, total):
    set_bg(slide, WHITE)
    add_header_bar(slide, "ZenseAI.QI – AI-Led QE Foundry", tag="PLATFORM",
                   subtitle="10+ AI agents — full Software Testing Lifecycle coverage")
    add_logo(slide)

    categories = [
        ("Requirements", "DeepSpeci", "Requirements Evaluator",
         "Scores & enhances user stories, surfaces ambiguity", CRIMSON),
        ("Test Design", "CaseGeni", "Test Case Generator",
         "Risk-based test design from AC & domain context", rgb("8B5CF6")),
        ("Automation", "Auto-PlayPilot", "Playwright Automation",
         "AI script gen, self-healing & MCP", rgb("059669")),
        ("Execution", "DataGeni", "Test Data Generator",
         "Synthetic data for edge & boundary cases", rgb("F59E0B")),
        ("Analysis", "Insights360", "Analytics & Reporting",
         "Release health dashboards & coverage analytics", rgb("EC4899")),
        ("Security", "Secure-Xi", "Security Agent",
         "OWASP-aligned security testing", rgb("EF4444")),
        ("Performance", "Perf-Xi", "Performance Agent",
         "AI script gen & execution analysis", rgb("F87171")),
    ]
    row2 = [
        ("Defect Intelligence", "Defect Intelligence",
         "Auto RCA & defect prediction", rgb("00BCD4")),
        ("Knowledge Base", "Knowledge Base",
         "Domain knowledge for accurate AI outcomes (RAG)", rgb("7C3AED")),
        ("RIA", "RIA", "Retrieval Intelligence Assistant — AI-powered QA colleague", GREEN),
    ]

    col_w = 12.6 / 7
    for ci, (cat, agent, sub, desc, clr) in enumerate(categories):
        cx = 0.35 + ci * col_w
        # Category header
        add_rect(slide, cx, 1.55, col_w - 0.05, 0.28, NAVY)
        add_text(slide, cat, cx + 0.05, 1.55, col_w - 0.1, 0.28,
                 font_size=8, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # Agent card
        add_rect(slide, cx, 1.88, col_w - 0.05, 1.85, CARD)
        add_text(slide, agent, cx + 0.07, 1.92, col_w - 0.12, 0.35,
                 font_size=13, bold=True, color=clr)
        add_text(slide, sub, cx + 0.07, 2.28, col_w - 0.12, 0.22,
                 font_size=8, color=MUTED)
        add_text(slide, desc, cx + 0.07, 2.52, col_w - 0.12, 1.15,
                 font_size=8.5, color=BODY, wrap=True)

    # Row 2
    r2_w = 12.6 / 3
    for ci, (cat, agent, desc, clr) in enumerate(row2):
        cx = 0.35 + ci * r2_w
        add_rect(slide, cx, 4.0, r2_w - 0.1, 0.28, NAVY)
        add_text(slide, cat, cx + 0.05, 4.0, r2_w - 0.1, 0.28,
                 font_size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_rect(slide, cx, 4.32, r2_w - 0.1, 1.5, CARD)
        add_text(slide, agent, cx + 0.1, 4.38, r2_w - 0.2, 0.35,
                 font_size=14, bold=True, color=clr)
        add_text(slide, desc, cx + 0.1, 4.78, r2_w - 0.2, 1.0,
                 font_size=9.5, color=BODY, wrap=True)

    add_footer(slide, slide_num, total)


def build_section_divider(slide, slide_num, total, tag: str, title: str, subtitle: str = ""):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, 5.8, NAVY)
    add_rect(slide, 0, 5.8, SLIDE_W, 0.08, CRIMSON)
    # Decorative circle
    c = slide.shapes.add_shape(9, Inches(9.0), Inches(-0.5), Inches(4.0), Inches(4.0))
    c.fill.solid(); c.fill.fore_color.rgb = rgb("002D5B"); c.line.fill.background()
    add_text(slide, tag, 0.9, 0.5, 4.0, 0.3, font_size=9, bold=True, color=CRIMSON)
    add_logo(slide)
    add_text(slide, title, 0.7, 2.3, 11.5, 1.8,
             font_size=44, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE, wrap=True)
    if subtitle:
        add_text(slide, subtitle, 0.7, 4.3, 11.5, 0.5,
                 font_size=16, color=rgb("B8C2E0"))
    add_footer(slide, slide_num, total)


def build_pov_plan(slide, slide_num, total):
    set_bg(slide, WHITE)
    add_header_bar(slide, "Proposed POV — Phased Implementation Plan", tag="PLATFORM",
                   subtitle="Building confidence incrementally; starting with zero-risk POC")
    add_logo(slide)

    phases = [
        {
            "phase": "Phase 1 — POC", "duration": "Weeks 1–4", "color": GREEN,
            "title": "PDF Validator + DeepSpeci Demo",
            "items": [
                "Deploy DocuProof on 50 Lace/Paragon docs",
                "DeepSpeci on BCAS epic — score AC gaps",
                "CaseGeni generates test cases from AC",
                "Zero infra cost — cloud sandbox",
            ],
            "kpi": "0 → 100% automated PDF checks"
        },
        {
            "phase": "Phase 2 — POV Expand", "duration": "Weeks 5–10", "color": NAVY,
            "title": "Auto-PlayPilot + Defect Intelligence",
            "items": [
                "Auto-PlayPilot on CAARS — Playwright MVP",
                "Defect Intelligence on Azure DevOps defects",
                "Self-healing runtime for BCAS/PayGW flakiness",
                "Knowledge Base seeded with tribal context",
            ],
            "kpi": "3 critical apps get automation baseline"
        },
        {
            "phase": "Phase 3 — Scale", "duration": "Weeks 11–16", "color": CRIMSON,
            "title": "Security, Performance & Insight360",
            "items": [
                "Secure-Xi on Payment Gateway + FHUB APIs",
                "Perf-Xi regression on high-traffic paths",
                "Insight360 replaces Excel sign-off dashboards",
                "RIA deployed as in-app QA assistant",
            ],
            "kpi": "Full STLC intelligence across 3 towers"
        },
    ]

    phase_w = 3.9
    gap = 0.27
    sx = 0.35
    sy = 1.55

    for i, ph in enumerate(phases):
        cx = sx + i * (phase_w + gap)
        clr = ph["color"]
        # Header
        add_rect(slide, cx, sy, phase_w, 0.55, clr)
        add_text(slide, ph["phase"], cx + 0.12, sy + 0.02, phase_w - 0.2, 0.28,
                 font_size=11, bold=True, color=WHITE)
        add_text(slide, ph["duration"], cx + 0.12, sy + 0.3, phase_w - 0.2, 0.22,
                 font_size=9, color=WHITE)
        # Card body
        card_h = 4.8
        add_rect(slide, cx, sy + 0.55, phase_w, card_h, CARD)
        add_text(slide, ph["title"], cx + 0.15, sy + 0.65, phase_w - 0.25, 0.4,
                 font_size=11, bold=True, color=NAVY, wrap=True)
        for j, item in enumerate(ph["items"]):
            add_text(slide, f"• {item}", cx + 0.15, sy + 1.12 + j * 0.5,
                     phase_w - 0.25, 0.45,
                     font_size=9.5, color=BODY, wrap=True)
        # KPI chip
        add_rect(slide, cx, sy + 0.55 + card_h - 0.32, phase_w, 0.32, clr)
        add_text(slide, ph["kpi"], cx + 0.1, sy + 0.55 + card_h - 0.32,
                 phase_w - 0.15, 0.32,
                 font_size=9, bold=True, color=WHITE,
                 valign=MSO_ANCHOR.MIDDLE)

    add_footer(slide, slide_num, total)


def build_pov_readiness(slide, slide_num, total):
    set_bg(slide, WHITE)
    add_header_bar(slide, "PoV Prerequisites, Skills & Success Measures", tag="PoV READINESS",
                   subtitle="What Zensar brings · What AA needs to provide · How we measure success")
    add_logo(slide)

    sections = [
        {
            "title": "What Zensar Brings", "color": NAVY,
            "items": [
                "ZenseAI.QI platform (cloud-hosted, zero AA infra needed for POC)",
                "2× QE Solution Architects + 1× AI Engineer",
                "Pre-built connectors: Jira, Zephyr, Azure DevOps, GitHub",
                "LLM layer (Azure OpenAI + Claude) — fully managed",
                "Template test strategies, onboarding runbooks",
            ]
        },
        {
            "title": "What AA Needs to Provide", "color": CRIMSON,
            "items": [
                "Jira project access (read) for BCAS epic",
                "Sample Lace/Paragon PDFs (50 docs for Phase 1)",
                "1× QE lead as point of contact (4h/week)",
                "InfoSec whitelist: Azure OpenAI endpoints",
                "Access to Azure DevOps pipelines for Phase 2",
            ]
        },
        {
            "title": "How We Measure Success", "color": GREEN,
            "items": [
                "PDF validation: 0 → 100% automated checks on 50 docs",
                "AC quality score baseline vs post-DeepSpeci",
                "CAARS: measurable automation coverage from 0%",
                "Flaky test reduction: BCAS pipeline green rate",
                "Qualitative: team confidence & time-to-test reduction",
            ]
        },
    ]

    col_w = 3.9
    gap = 0.27
    sx = 0.35

    for i, sec in enumerate(sections):
        cx = sx + i * (col_w + gap)
        add_rect(slide, cx, 1.55, col_w, 0.38, sec["color"])
        add_text(slide, sec["title"], cx + 0.12, 1.55, col_w - 0.2, 0.38,
                 font_size=12, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
        card_h = 4.85
        add_rect(slide, cx, 1.93, col_w, card_h, CARD)
        for j, item in enumerate(sec["items"]):
            add_text(slide, f"• {item}", cx + 0.15, 2.03 + j * 0.88,
                     col_w - 0.25, 0.82, font_size=10, color=BODY, wrap=True)

    add_footer(slide, slide_num, total)


def build_closing(slide, slide_num, total):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, 5.8, NAVY)
    add_rect(slide, 0, 5.8, SLIDE_W, 0.08, CRIMSON)
    c = slide.shapes.add_shape(9, Inches(9.0), Inches(-0.5), Inches(4.5), Inches(4.5))
    c.fill.solid(); c.fill.fore_color.rgb = rgb("002D5B"); c.line.fill.background()
    add_text(slide, "ZenseAI.QI  •  PROPOSAL", 0.9, 0.5, 4.0, 0.3,
             font_size=9, bold=True, color=CRIMSON)
    add_logo(slide)
    add_text(slide, "Thank you", 0.7, 2.0, 11.5, 1.5,
             font_size=52, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "JUNE 2026", 0.9, 6.0, 3.0, 0.3,
             font_size=11, bold=True, color=CRIMSON)
    add_footer(slide, slide_num, total)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    TOTAL = 12  # total slides we'll generate

    # 1. Cover
    s = blank_slide(prs); build_cover(s, 1, TOTAL)

    # 2. Agenda
    s = blank_slide(prs); build_agenda(s, 2, TOTAL)

    # 3. AS-IS pipeline — Testing Lifecycle Today
    s = blank_slide(prs)
    build_pipeline_slide(s, 3, TOTAL,
        tag="AS-IS  THE AA",
        title="Testing Lifecycle Today at The AA",
        cols_data=[
            {"agent": "Jira + manual analysis", "badge": "Manual", "badge_color": "F59E0B",
             "points": ["Stories authored in Jira", "RoVo agents piloted", "Acceptance criteria written ad-hoc", "No automated story → TC traceability"]},
            {"agent": "Zephyr + Excel", "badge": "Manual", "badge_color": "F59E0B",
             "points": ["TCs typed in Zephyr", "Duplicated Excel sign-off", "CAARS: 500+ TCs manual", "Payment Gateway: 6000+ TCs"]},
            {"agent": "Playwright + C# (.NET)", "badge": "Strong",
             "points": ["GHCP for Selenium → Playwright", "Vehicle Insp.: 80–85% automated", "BCAS, Payment Gateway, FHUB automated"]},
            {"agent": "CI/CD (Azure DevOps)", "badge": "Daily CI",
             "points": ["Daily pipelines for Veh Insp/FHUB/PayGW", "Flaky failures: BCAS, Payment GW", "Manual execution for CAARS, Lace, CAIS", "No mobile execution"]},
            {"agent": "Jira + Azure DevOps", "badge": "Reactive", "badge_color": "F59E0B",
             "points": ["Defects raised manually", "No prediction / risk-based triage", "Rerun cost high on flaky pipelines", "Limited cross-portfolio defect insight"]},
            {"agent": "Manual + Applitools", "badge": "Manual", "badge_color": "F59E0B",
             "points": ["PDF Validation: Applitools (manual Lace 500+)", "SAP: exploring TOSCA + Cloud ALM", "No Security / Perf in scope", "Reporting: Excel + Zephyr duplication"]},
        ],
        bottom_row=[("8 testers · 3 portfolios", "Enterprise · Insurance · Aurora"),
                    ("3 critical apps with 0% automation", "CAARS · Lace · CAIS")])

    # 4. Pain points
    s = blank_slide(prs); build_challenges(s, 4, TOTAL)

    # 5. TO-BE pipeline — ZenseAI.QI agents
    s = blank_slide(prs)
    build_pipeline_slide(s, 5, TOTAL,
        tag="TO-BE  THE AA",
        title="Testing Lifecycle Reimagined with AA-AI.QI",
        cols_data=[
            {"agent": "DeepSpeci", "points": ["Stories ↔ AC w/ traceability", "Auto-extracts AC from Jira", "Feeds CaseGeni downstream"]},
            {"agent": "CaseGeni", "points": ["Structured AI test-case generation", "Traceable test cases", "Risk-based coverage"]},
            {"agent": "Auto-PlayPilot", "points": ["Playwright + C# scaffolding", "Day-1 framework", "GHCP-friendly output"]},
            {"agent": "Self-Healing Exec", "points": ["Auto-PlayPilot runtime", "Self-healing locators", "Cuts rerun cost on flaky suites"]},
            {"agent": "Defect Intelligence", "points": ["Triage + risk prediction", "Predictive risk on changed code", "Auto RCA"]},
            {"agent": "DocuProof · Secure-Xi · Perf-Xi", "points": ["Visual+LLM diff for PDFs", "Security Agents (OWASP)", "Performance Agent", "Sec/Perf zero-scope value-adds"]},
        ],
        bottom_row=[("Knowledge Base (RAG)", "Captures tribal knowledge — SAP, Lace, CAARS"),
                    ("Virtual In-app assistant", "Live answers + failure explanations to testers"),
                    ("Insight360 dashboard", "Cross-tower QI observability — replaces Excel sign-off")])

    # 6. Platform overview
    s = blank_slide(prs); build_platform_overview(s, 6, TOTAL)

    # 7. Demo section divider
    s = blank_slide(prs)
    build_section_divider(s, 7, TOTAL,
                          tag="ZenseAI.QI  •  PROPOSAL",
                          title="Demo\nZenseAI.QI Foundry",
                          subtitle="JUNE 2026")

    # 8. Next Steps pipeline (POV roadmap per agent)
    s = blank_slide(prs)
    build_pipeline_slide(s, 8, TOTAL,
        tag="TO-BE  THE AA",
        title="Next Steps – POV – AA.AI-QI",
        cols_data=[
            {"agent": "DeepSpeci", "badge": "POV",
             "points": ["Stories ↔ AC w/ traceability", "Auto-extracts AC from Jira", "Feeds CaseGeni downstream"],
             "poc_target": "BCAS · Cathie · Aurora epic",
             "later": "Pull tribal SAP context via KB (H3)"},
            {"agent": "CaseGeni", "badge": "POV",
             "points": ["Structured AI test-case generation", "Traceable test cases"],
             "poc_target": "Same epic — chained pipeline",
             "later": "Roll out across all 3 towers"},
            {"agent": "Auto-PlayPilot", "badge": "POV",
             "points": ["Playwright + C# scaffolding", "Day-1 framework", "GHCP-friendly output"],
             "poc_target": "CAARS automation MVP",
             "later": "Mobile capability ramp (H3)"},
            {"agent": "Self-Healing Exec", "badge": "Later — H2", "badge_color": "F59E0B",
             "points": ["Auto-PlayPilot runtime", "Self-healing locators", "Cuts rerun cost on flaky suites"],
             "later": "Target: BCAS + Payment Gateway"},
            {"agent": "Defect Intelligence", "badge": "Later — H2", "badge_color": "F59E0B",
             "points": ["Triage + risk prediction", "Predictive risk on changed code"],
             "later": "All 3 towers post POV"},
            {"agent": "DocuProof / Secure-Xi", "badge": "POV",
             "points": ["PDF visual+LLM diff on Lace docs", "OWASP security scan on PayGW"],
             "poc_target": "50 Lace docs Phase 1",
             "later": "Expand to full PDF suite"},
        ])

    # 9. Phased POV plan
    s = blank_slide(prs); build_pov_plan(s, 9, TOTAL)

    # 10. PoV readiness
    s = blank_slide(prs); build_pov_readiness(s, 10, TOTAL)

    # 11. Infra requirements — simple table slide
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_header_bar(s, "Infrastructure & Access Requirements", tag="INFRA",
                   subtitle="What needs to be in place before kick-off")
    add_logo(s)
    infra_rows = [
        ("Azure OpenAI", "Endpoint + API key", "InfoSec whitelist required", "Phase 1"),
        ("Jira / Zephyr", "Read access to BCAS project", "Service account", "Phase 1"),
        ("Azure DevOps", "Pipeline read + webhook", "PAT token", "Phase 2"),
        ("GitHub", "Repo read (for Playwright scripts)", "PAT or GitHub App", "Phase 2"),
        ("Network", "Outbound HTTPS to Azure AI endpoints", "Proxy/firewall exception", "Phase 1"),
        ("VMs / Agents", "1× Azure DevOps build agent (our image)", "Zensar provisions", "Phase 1"),
    ]
    headers = ["Component", "Requirement", "Access Type", "Phase"]
    hw = [3.2, 3.8, 2.8, 1.5]
    hx = 0.35
    hy = 1.6
    for ci, (hdr, w) in enumerate(zip(headers, hw)):
        add_rect(s, hx, hy, w - 0.05, 0.35, NAVY)
        add_text(s, hdr, hx + 0.1, hy, w - 0.15, 0.35,
                 font_size=10, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
        hx += w
    for ri, row in enumerate(infra_rows):
        ry = hy + 0.35 + ri * 0.46
        rx = 0.35
        fill = CARD if ri % 2 == 0 else WHITE
        for ci, (cell, w) in enumerate(zip(row, hw)):
            add_rect(s, rx, ry, w - 0.05, 0.42, fill)
            add_text(s, cell, rx + 0.1, ry + 0.02, w - 0.18, 0.38,
                     font_size=9.5, color=BODY, valign=MSO_ANCHOR.MIDDLE)
            rx += w
    add_footer(s, 11, TOTAL)

    # 12. Closing
    s = blank_slide(prs); build_closing(s, 12, TOTAL)

    out = "/tmp/AA_ZenseAI_QI_Zensar_Template.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
