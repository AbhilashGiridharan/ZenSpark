from pptx import Presentation
from pptx.util import Inches, Pt
import json

src = "/Users/AG125765/Library/CloudStorage/OneDrive-ZensarTechnologiesLtd/Zensar Docs/Proposals/2026/AA/AA_ZenseAI.QI Offerings V7.0.pptx"
prs = Presentation(src)
slides = list(prs.slides)

result = []
for i, slide in enumerate(slides):
    slide_data = {"slide": i+1, "shapes": []}
    for shape in slide.shapes:
        stype = shape.shape_type
        name = shape.name
        try:
            l = round(shape.left.inches, 3) if shape.left else None
            t = round(shape.top.inches, 3) if shape.top else None
            w = round(shape.width.inches, 3) if shape.width else None
            h = round(shape.height.inches, 3) if shape.height else None
        except:
            l = t = w = h = None

        s = {"type": stype, "name": name, "pos": [l, t, w, h]}

        try:
            f = shape.fill
            if f.type is not None:
                try:
                    s["fill"] = "#" + str(f.fore_color.rgb)
                except:
                    pass
        except:
            pass

        if hasattr(shape, "text_frame") and shape.text_frame:
            paras = []
            for para in shape.text_frame.paragraphs:
                if not para.text.strip():
                    continue
                p_info = {"text": para.text}
                for run in para.runs[:1]:
                    try:
                        if run.font.size:
                            p_info["size"] = round(run.font.size.pt, 1)
                        if run.font.bold:
                            p_info["bold"] = True
                        try:
                            p_info["color"] = "#" + str(run.font.color.rgb)
                        except:
                            pass
                    except:
                        pass
                paras.append(p_info)
            if paras:
                s["text_paras"] = paras

        if stype == 13:
            s["is_image"] = True

        if "text_paras" in s or stype == 13:
            slide_data["shapes"].append(s)

    result.append(slide_data)

with open("/tmp/src_content.json", "w") as f:
    json.dump(result, f, indent=2)
print("Done. Written to /tmp/src_content.json")
print(f"Total slides: {len(result)}")
for r in result:
    texts = [p["text"][:60] for s in r["shapes"] if "text_paras" in s for p in s["text_paras"][:2]]
    print(f"  Slide {r['slide']}: {texts[:4]}")
